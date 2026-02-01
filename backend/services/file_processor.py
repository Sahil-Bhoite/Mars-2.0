"""
File Processing Service - Extract text from various document types.
Based on M.A.R.S original implementation by Sahil Bhoite.

Supports 35+ file types including documents, code, audio, video, and images.
"""
import io
import json
import zipfile
from typing import Optional

from config import get_settings

settings = get_settings()


def simple_chunk_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list[str]:
    """Simple text chunking without langchain dependency."""
    if not text:
        return []
    
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        start = end - chunk_overlap
        if start < 0:
            start = 0
    return chunks


class FileProcessor:
    """Process and extract text from various file types."""
    
    # All supported file extensions - 35+ types
    SUPPORTED_EXTENSIONS = {
        # Documents
        "pdf", "docx", "doc", "pptx", "ppt", "xlsx", "xls", "csv", "txt",
        # Code files - Extended
        "py", "java", "c", "h", "cpp", "js", "ts", "swift", "r", "rs", "sql",
        "go", "kt", "scala", "php", "rb", "sh", "bash",
        # Web/Markup
        "html", "css", "json", "xml", "md", "tex",
        # Images (OCR)
        "jpg", "jpeg", "png", "bmp",
        # Archives
        "zip",
        # Audio
        "wav", "mp3",
        # Video
        "mp4"
    }
    
    def __init__(self):
        self.chunk_size = settings.chunk_size
        self.chunk_overlap = settings.chunk_overlap
    
    def extract_text(self, file_content: bytes, filename: str) -> str:
        """Extract text from a file based on its extension."""
        extension = filename.split(".")[-1].lower()
        
        try:
            # Documents
            if extension == "pdf":
                return self._extract_pdf(file_content)
            elif extension == "docx":
                return self._extract_docx(file_content)
            elif extension == "doc":
                return self._extract_doc_legacy(file_content)
            elif extension == "pptx":
                return self._extract_pptx(file_content)
            elif extension == "ppt":
                return self._extract_ppt_legacy(file_content)
            elif extension in ["xlsx", "xls"]:
                return self._extract_excel(file_content)
            elif extension == "csv":
                return self._extract_csv(file_content)
            
            # Code files - plain text (extended)
            elif extension in ["py", "java", "c", "h", "cpp", "js", "ts", "swift", "r", "rs", "sql", 
                             "css", "txt", "go", "kt", "scala", "php", "rb", "sh", "bash"]:
                return self._extract_code(file_content, extension)
            
            # Markup
            elif extension == "html":
                return self._extract_html(file_content)
            elif extension == "xml":
                return self._extract_xml(file_content)
            elif extension == "json":
                return self._extract_json(file_content)
            elif extension in ["md", "tex"]:
                return self._extract_text_plain(file_content)
            
            # Images with OCR
            elif extension in ["jpg", "jpeg", "png", "bmp"]:
                return self._extract_image_ocr(file_content)
            
            # Archives
            elif extension == "zip":
                return self._extract_zip(file_content)
            
            # Audio (WAV and MP3)
            elif extension in ["wav", "mp3"]:
                return self._extract_audio(file_content, extension)
            
            # Video (MP4)
            elif extension == "mp4":
                return self._extract_video(file_content)
            
            else:
                raise ValueError(f"Unsupported file type: {extension}")
                
        except Exception as e:
            raise ValueError(f"Error processing {filename}: {str(e)}")
    
    def _extract_pdf(self, content: bytes) -> str:
        """Extract text from PDF file."""
        try:
            from PyPDF2 import PdfReader
            text = ""
            pdf_reader = PdfReader(io.BytesIO(content))
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            if not text.strip():
                return "[PDF file - no extractable text found]"
            return text.strip()
        except Exception as e:
            raise ValueError(f"PDF processing error: {str(e)}")
    
    def _extract_docx(self, content: bytes) -> str:
        """Extract text from Word document."""
        try:
            from docx import Document
            text = ""
            doc = Document(io.BytesIO(content))
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    text += row_text + "\n"
            return text.strip() if text.strip() else "[DOCX file - no text found]"
        except Exception as e:
            raise ValueError(f"DOCX processing error: {str(e)}")
    
    def _extract_pptx(self, content: bytes) -> str:
        """Extract text from PowerPoint presentation."""
        try:
            from pptx import Presentation
            text = ""
            prs = Presentation(io.BytesIO(content))
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n--- Slide {slide_num} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            return text.strip() if text.strip() else "[PPTX file - no text found]"
        except Exception as e:
            raise ValueError(f"PPTX processing error: {str(e)}")
    
    def _extract_doc_legacy(self, content: bytes) -> str:
        """Extract text from legacy .doc file (Word 97-2003)."""
        try:
            import subprocess
            import tempfile
            import os
            
            # Write to temp file
            with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as f:
                f.write(content)
                temp_path = f.name
            
            try:
                # Try antiword first (if installed)
                result = subprocess.run(['antiword', temp_path], capture_output=True, text=True, timeout=30)
                if result.returncode == 0 and result.stdout.strip():
                    return result.stdout.strip()
            except (FileNotFoundError, subprocess.TimeoutExpired):
                pass
            
            try:
                # Try catdoc as fallback
                result = subprocess.run(['catdoc', temp_path], capture_output=True, text=True, timeout=30)
                if result.returncode == 0 and result.stdout.strip():
                    return result.stdout.strip()
            except (FileNotFoundError, subprocess.TimeoutExpired):
                pass
            
            # Last resort: try to read as binary and extract text
            try:
                text = content.decode('utf-8', errors='ignore')
                # Filter printable characters
                printable = ''.join(c for c in text if c.isprintable() or c in '\n\t ')
                if len(printable) > 100:
                    return f"[Legacy DOC - partial extraction]\n{printable[:5000]}"
            except:
                pass
            
            os.unlink(temp_path)
            return "[Legacy .doc file - install 'antiword' or 'catdoc' for full extraction]"
        except Exception as e:
            return f"[DOC processing limited: {str(e)}]"
    
    def _extract_ppt_legacy(self, content: bytes) -> str:
        """Extract text from legacy .ppt file (PowerPoint 97-2003)."""
        try:
            # Try to extract any readable text from binary
            text = content.decode('utf-8', errors='ignore')
            printable = ''.join(c for c in text if c.isprintable() or c in '\n\t ')
            
            # Filter very short strings
            words = [w for w in printable.split() if len(w) > 2]
            if len(words) > 20:
                return f"[Legacy PPT - partial extraction]\n{' '.join(words[:500])}"
            
            return "[Legacy .ppt file - convert to .pptx for full extraction]"
        except Exception as e:
            return f"[PPT processing limited: {str(e)}]"
    
    def _extract_excel(self, content: bytes) -> str:
        """Extract text from Excel file."""
        try:
            import pandas as pd
            text = ""
            excel_file = io.BytesIO(content)
            df_dict = pd.read_excel(excel_file, sheet_name=None, engine='openpyxl')
            for sheet_name, df in df_dict.items():
                text += f"\n--- Sheet: {sheet_name} ---\n"
                text += df.to_string(index=False) + "\n"
            return text.strip() if text.strip() else "[Excel file - no data found]"
        except Exception as e:
            raise ValueError(f"Excel processing error: {str(e)}")
    
    def _extract_csv(self, content: bytes) -> str:
        """Extract text from CSV file."""
        try:
            import pandas as pd
            # Try different encodings
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    df = pd.read_csv(io.BytesIO(content), encoding=encoding)
                    return df.to_string(index=False)
                except UnicodeDecodeError:
                    continue
            # Fallback: just read as text
            return content.decode('latin-1', errors='ignore')
        except Exception as e:
            raise ValueError(f"CSV processing error: {str(e)}")
    
    def _extract_code(self, content: bytes, extension: str) -> str:
        """Extract code with language context."""
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("latin-1", errors='ignore')
        
        lang_names = {
            "py": "Python", "java": "Java", "c": "C", "h": "C Header",
            "cpp": "C++", "js": "JavaScript", "ts": "TypeScript",
            "swift": "Swift", "r": "R", "rs": "Rust", "sql": "SQL",
            "css": "CSS", "txt": "Text",
            # New languages
            "go": "Go", "kt": "Kotlin", "scala": "Scala",
            "php": "PHP", "rb": "Ruby", "sh": "Shell", "bash": "Bash"
        }
        
        lang = lang_names.get(extension, extension.upper())
        return f"[{lang} Code]\n{text}"
    
    def _extract_html(self, content: bytes) -> str:
        """Extract text from HTML file."""
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(content, 'html.parser')
            # Remove script and style elements
            for element in soup(['script', 'style']):
                element.decompose()
            return soup.get_text(separator='\n', strip=True)
        except Exception as e:
            # Fallback to plain text
            return self._extract_text_plain(content)
    
    def _extract_xml(self, content: bytes) -> str:
        """Extract text from XML file."""
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(content, 'lxml-xml')
            return soup.get_text(separator='\n', strip=True)
        except Exception as e:
            # Fallback to plain text
            return self._extract_text_plain(content)
    
    def _extract_json(self, content: bytes) -> str:
        """Extract and format JSON file."""
        try:
            text = content.decode("utf-8")
            data = json.loads(text)
            return json.dumps(data, indent=2)
        except Exception as e:
            return self._extract_text_plain(content)
    
    def _extract_text_plain(self, content: bytes) -> str:
        """Extract plain text."""
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            return content.decode("latin-1", errors='ignore')
    
    def _extract_image_ocr(self, content: bytes) -> str:
        """Extract text from image using OCR."""
        try:
            import cv2
            import pytesseract
            import numpy as np
            
            nparr = np.frombuffer(content, np.uint8)
            image = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            text = pytesseract.image_to_string(gray)
            return f"[OCR Extracted Text]\n{text}" if text.strip() else "[Image - no text extracted]"
        except ImportError:
            return "[OCR not available - install opencv-python and pytesseract]"
        except Exception as e:
            return f"[Image processing error: {str(e)}]"
    
    def _extract_zip(self, content: bytes) -> str:
        """Extract text from files inside ZIP archive."""
        text = ""
        try:
            with zipfile.ZipFile(io.BytesIO(content), 'r') as z:
                for filename in z.namelist():
                    if filename.endswith('/'):
                        continue
                    ext = filename.split('.')[-1].lower()
                    if ext in self.SUPPORTED_EXTENSIONS and ext != 'zip':
                        try:
                            with z.open(filename) as f:
                                file_content = f.read()
                                extracted = self.extract_text(file_content, filename)
                                text += f"\n--- {filename} ---\n{extracted}\n"
                        except:
                            pass
            return text.strip() if text.strip() else "[ZIP file - no extractable content]"
        except Exception as e:
            raise ValueError(f"ZIP processing error: {str(e)}")
    
    def _extract_audio(self, content: bytes, extension: str) -> str:
        """Extract text from audio (WAV/MP3) using speech recognition."""
        try:
            import speech_recognition as sr
            import tempfile
            import os
            
            # For MP3, we need to convert to WAV first
            if extension == "mp3":
                try:
                    from pydub import AudioSegment
                    audio = AudioSegment.from_mp3(io.BytesIO(content))
                    wav_io = io.BytesIO()
                    audio.export(wav_io, format="wav")
                    content = wav_io.getvalue()
                except ImportError:
                    return "[MP3 processing requires pydub - install with: pip install pydub]"
            
            with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            
            try:
                r = sr.Recognizer()
                with sr.AudioFile(tmp_path) as source:
                    audio_data = r.record(source)
                    text = r.recognize_google(audio_data)
                return f"[Audio Transcription]\n{text}"
            finally:
                os.unlink(tmp_path)
        except ImportError:
            return "[Audio processing not available - install speech_recognition]"
        except Exception as e:
            return f"[Audio processing error: {str(e)}]"
    
    def _extract_video(self, content: bytes) -> str:
        """Extract audio from video and transcribe."""
        try:
            from moviepy.editor import VideoFileClip
            import tempfile
            import os
            
            # Save video to temp file
            with tempfile.NamedTemporaryFile(suffix='.mp4', delete=False) as tmp_video:
                tmp_video.write(content)
                tmp_video_path = tmp_video.name
            
            try:
                # Extract audio
                video = VideoFileClip(tmp_video_path)
                
                with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as tmp_audio:
                    tmp_audio_path = tmp_audio.name
                
                video.audio.write_audiofile(tmp_audio_path, verbose=False, logger=None)
                video.close()
                
                # Transcribe audio
                with open(tmp_audio_path, 'rb') as f:
                    audio_content = f.read()
                
                result = self._extract_audio(audio_content, 'wav')
                os.unlink(tmp_audio_path)
                
                return f"[Video Transcription]\n{result.replace('[Audio Transcription]', '').strip()}"
            finally:
                os.unlink(tmp_video_path)
        except ImportError:
            return "[Video processing requires moviepy - install with: pip install moviepy]"
        except Exception as e:
            return f"[Video processing error: {str(e)}]"
    
    def chunk_text(self, text: str, filename: str) -> list[dict]:
        """Split text into chunks with metadata."""
        if not text or not text.strip():
            return []
        chunks = simple_chunk_text(text, self.chunk_size, self.chunk_overlap)
        return [
            {
                "text": chunk,
                "source": filename,
                "chunk_index": i
            }
            for i, chunk in enumerate(chunks)
        ]


# Singleton instance
file_processor = FileProcessor()
